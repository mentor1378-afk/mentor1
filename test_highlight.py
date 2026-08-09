from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

NS_A14 = 'http://schemas.microsoft.com/office/drawing/2010/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def add_highlight(run, color_hex):
    """Use solidFill on the run's rPr as a background color workaround for PowerPoint"""
    rPr = run._r.get_or_add_rPr()
    
    # Remove existing solidFill if any (used as highlight)
    for sf in rPr.findall(f'{{{NS_A}}}solidFill'):
        rPr.remove(sf)
    
    # For PowerPoint, we set font color explicitly and use a different method
    # PowerPoint doesn't support a:highlight natively on slides
    # Instead we'll use the run's font color for text and rely on shape-level highlighting
    # 
    # Actually, let's try the MC AlternateContent approach with a14:textFill
    pass

def add_highlight_via_shape(slide, left, top, width, height, color_hex):
    """Create a colored rectangle behind text as a highlight"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()  # No border
    return shape

# Test: Create colored text runs with font color instead of highlight
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

from pptx.enum.shapes import MSO_SHAPE

# Gray content box
body_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5))
body_box.fill.solid()
body_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
body_box.line.color.rgb = RGBColor(200, 200, 200)

# Green highlight bar behind "Welcome" text
hl1 = add_highlight_via_shape(slide, Inches(0.6), Inches(0.6), Inches(6.5), Inches(0.45), '00FF00')

# Text box on top
tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(6))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Welcome to SIMATS ENGINEERING"
run.font.bold = True
run.font.size = Pt(22)
run.font.name = "Times New Roman"

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Dear Parent,"
run2.font.size = Pt(18)
run2.font.bold = True
run2.font.name = "Times New Roman"

# Green highlight for attendance line
hl2 = add_highlight_via_shape(slide, Inches(0.6), Inches(1.55), Inches(10.0), Inches(0.4), '00FF00')

p3 = tf.add_paragraph()
run3 = p3.add_run()
run3.text = "So far the student has maintained consistent attendance in the course."
run3.font.size = Pt(18)
run3.font.bold = True
run3.font.name = "Times New Roman"

prs.save('test_shape_highlight.pptx')
print("Saved test_shape_highlight.pptx")
print("This uses colored shapes behind the text as PowerPoint-native highlights.")
