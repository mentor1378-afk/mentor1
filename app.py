import os
import json
import urllib.request
from io import BytesIO
from flask import Flask, render_template, request, redirect, url_for, session, flash, send_file
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from werkzeug.utils import secure_filename
import cloudinary
import cloudinary.uploader
import cloudinary.api
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from datetime import datetime

def add_highlight(run, color_hex):
    rPr = run._r.get_or_add_rPr()
    for hl in rPr.xpath('./a:highlight'):
        rPr.remove(hl)
        
    highlight = OxmlElement('a:highlight')
    srgbClr = OxmlElement('a:srgbClr')
    srgbClr.set('val', color_hex)
    highlight.append(srgbClr)
    
    tags_after = [
        '{http://schemas.openxmlformats.org/drawingml/2006/main}uLnTx',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}latin',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}ea',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}cs',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}sym',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkMouseOver',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}rtl',
        '{http://schemas.openxmlformats.org/drawingml/2006/main}extLst',
    ]
    
    inserted = False
    for child in rPr:
        if child.tag in tags_after:
            child.addprevious(highlight)
            inserted = True
            break
            
    if not inserted:
        rPr.append(highlight)

MENTOR_NAMES = {
    'mentor1@123': 'vijayavaran',
    'mentor2@123': 'Mr. Mentor 2'
}

app = Flask(__name__)
app.config['SECRET_KEY'] = 'cinematic_secret_key_123'
uri = os.environ.get("DATABASE_URL", "sqlite:///database.db")
if uri.startswith("postgres://"):
    uri = uri.replace("postgres://", "postgresql://", 1)
app.config['SQLALCHEMY_DATABASE_URI'] = uri
app.config['SQLALCHEMY_ENGINE_OPTIONS'] = {
    "pool_pre_ping": True,
    "pool_recycle": 300
}
app.config['UPLOAD_FOLDER'] = 'static/uploads'
app.config['ALLOWED_EXTENSIONS'] = {'png', 'jpg', 'jpeg', 'gif'}
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

db = SQLAlchemy(app)
bcrypt = Bcrypt(app)

# Models
class User(db.Model):
    username = db.Column(db.String(50), primary_key=True)
    password_hash = db.Column(db.String(100), nullable=False)
    role = db.Column(db.String(20), nullable=False) # 'faculty' or 'student'

class StudentDetail(db.Model):
    reg_num = db.Column(db.String(50), primary_key=True)
    name = db.Column(db.String(100))
    course = db.Column(db.String(100))
    mentor_username = db.Column(db.String(50), db.ForeignKey('user.username'))

    # JSON strings for dynamic data
    slot_info = db.Column(db.Text, default='[]') # List of slot names
    attendance_data = db.Column(db.Text, default='{}') # {"Slot A": 81}
    marks_data = db.Column(db.Text, default='{}') # {"Slot A": {"model": 20}}
    
    registered_new_course = db.Column(db.String(200))
    online_course = db.Column(db.String(200))
    event_participation = db.Column(db.Text)
    additional_description = db.Column(db.Text)
    photo_path = db.Column(db.String(200))

    def get_attendance(self):
        try: return json.loads(self.attendance_data)
        except: return {}
    
    def get_marks(self):
        try: return json.loads(self.marks_data)
        except: return {}
    
    def get_slots(self):
        try: return json.loads(self.slot_info)
        except: return []

# Helper functions
def delete_photo(photo_path):
    if not photo_path:
        return
    if 'cloudinary.com' in photo_path:
        try:
            public_id = "simats_profiles/" + photo_path.split('/')[-1].split('.')[0]
            cloudinary.uploader.destroy(public_id)
        except Exception as e:
            print(f'Failed to delete from Cloudinary: {e}')
    elif os.path.exists(photo_path):
        try:
            os.unlink(photo_path)
        except Exception as e:
            print(f'Failed to delete local photo: {e}')

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def seed_db():
    # Faculty: mentor1@123
    if not User.query.get('mentor1@123'):
        db.session.add(User(
            username='mentor1@123',
            password_hash=bcrypt.generate_password_hash('welcome').decode('utf-8'),
            role='faculty'
        ))

    # Faculty: mentor2@123
    if not User.query.get('mentor2@123'):
        db.session.add(User(
            username='mentor2@123',
            password_hash=bcrypt.generate_password_hash('welcome').decode('utf-8'),
            role='faculty'
        ))



    db.session.commit()

with app.app_context():
    db.create_all()  # Only creates tables if they don't exist — never wipes data
    seed_db()

# Routes
@app.route('/')
def index():
    if 'user' in session:
        if session['role'] == 'faculty':
            return redirect(url_for('faculty_dashboard'))
        return redirect(url_for('student_dashboard'))
    return redirect(url_for('login'))

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        user = User.query.filter_by(username=username).first()
        
        if user and bcrypt.check_password_hash(user.password_hash, password):
            session['user'] = user.username
            session['role'] = user.role
            if user.role == 'faculty':
                return redirect(url_for('faculty_dashboard'))
            return redirect(url_for('student_dashboard'))
        flash('Invalid username or password', 'danger')
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))

@app.route('/change_password', methods=['POST'])
def change_password():
    if 'user' not in session:
        return redirect(url_for('login'))
    
    current_password = request.form['current_password']
    new_password = request.form['new_password']
    
    user = User.query.get(session['user'])
    if user and bcrypt.check_password_hash(user.password_hash, current_password):
        user.password_hash = bcrypt.generate_password_hash(new_password).decode('utf-8')
        db.session.commit()
        flash('Password changed successfully!', 'success')
    else:
        flash('Incorrect current password', 'danger')
    
    return redirect(url_for('faculty_dashboard' if session['role'] == 'faculty' else 'student_dashboard'))

@app.route('/student', methods=['GET', 'POST'])
def student_dashboard():
    if 'user' not in session or session['role'] != 'student':
        return redirect(url_for('login'))
    
    student = StudentDetail.query.get(session['user'])
    if not student:
        student = StudentDetail(reg_num=session['user'], name=session['user'])
        db.session.add(student)
        db.session.commit()

    if request.method == 'POST':
        student.name = request.form.get('name', student.name) or student.name
        student.course = request.form.get('course', student.course) or student.course

        # Dynamic Slots handling
        slots = request.form.getlist('slot_names[]')
        if not slots:
            slots = student.get_slots() or []
        student.slot_info = json.dumps(slots)

        att_data = {}
        marks_data = {}
        current_attendance = student.get_attendance()
        current_marks = student.get_marks()
        for slot in slots:
            if not slot:
                continue
            att_data[slot] = request.form.get(f'att_{slot}', current_attendance.get(slot, 0))
            marks_data[slot] = {
                'total_marks': request.form.get(f'total_marks_{slot}', current_marks.get(slot, {}).get('total_marks', '-')),
                'model': request.form.get(f'model_{slot}', current_marks.get(slot, {}).get('model', '-')),
                'test1': request.form.get(f'test1_{slot}', current_marks.get(slot, {}).get('test1', '-')),
                'test2': request.form.get(f'test2_{slot}', current_marks.get(slot, {}).get('test2', '-')),
                'avg': request.form.get(f'avg_{slot}', current_marks.get(slot, {}).get('avg', '-')),
                'course': request.form.get(f'course_{slot}', current_marks.get(slot, {}).get('course', '-'))
            }

        student.attendance_data = json.dumps(att_data)
        student.marks_data = json.dumps(marks_data)

        student.registered_new_course = request.form.get('registered_new_course', student.registered_new_course or '')
        student.online_course = request.form.get('online_course', student.online_course or '')
        student.event_participation = request.form.get('event_participation', student.event_participation or '')
        student.additional_description = request.form.get('description', student.additional_description or '')

        file = request.files.get('photo')
        if file and file.filename and allowed_file(file.filename):
            if os.environ.get('CLOUDINARY_URL'):
                try:
                    upload_result = cloudinary.uploader.upload(file, folder="simats_profiles")
                    student.photo_path = upload_result.get('secure_url')
                except Exception as e:
                    flash(f'Failed to upload to Cloudinary. Please check your CLOUDINARY_URL. Error: {str(e)}', 'danger')
            else:
                filename = secure_filename(f"{session['user']}_{file.filename}")
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                student.photo_path = filepath

        db.session.commit()
        flash('Details updated successfully!', 'success')
        return redirect(url_for('student_dashboard'))
    
    # Passing current data as dicts
    return render_template('student.html', student=student, 
                           attendance=student.get_attendance(), 
                           marks=student.get_marks(), 
                           slots=student.get_slots())

@app.route('/add_student', methods=['POST'])
def add_student():
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    reg_num = request.form['reg_num']
    name = request.form['name']
    password = request.form['password']
    
    if User.query.get(reg_num):
        flash('Student with this Registration Number already exists.', 'danger')
    else:
        new_user = User(
            username=reg_num,
            password_hash=bcrypt.generate_password_hash(password).decode('utf-8'),
            role='student'
        )
        new_detail = StudentDetail(reg_num=reg_num, name=name, mentor_username=session['user'])
        db.session.add(new_user)
        db.session.add(new_detail)
        db.session.commit()
        flash(f'Student {name} added successfully!', 'success')
    
    return redirect(url_for('faculty_dashboard'))

@app.route('/edit_student_faculty/<reg_num>', methods=['POST'])
def edit_student_faculty(reg_num):
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    student = StudentDetail.query.get(reg_num)
    if student:
        student.name = request.form.get('name', student.name) or student.name
        student.course = request.form.get('course', student.course) or student.course
        student.registered_new_course = request.form.get('registered_new_course', student.registered_new_course or '')
        student.online_course = request.form.get('online_course', student.online_course or '')
        student.event_participation = request.form.get('event_participation', student.event_participation or '')
        student.additional_description = request.form.get('description', student.additional_description or '')
        
        db.session.commit()
        flash(f'Details for {student.name or reg_num} updated successfully!', 'success')
    else:
        flash('Student not found.', 'danger')
        
    return redirect(url_for('faculty_dashboard'))

@app.route('/faculty')
def faculty_dashboard():
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    students = StudentDetail.query.filter_by(mentor_username=session['user']).all()
    # Simplified stats for global update
    total_att_a = 0
    count_a = 0
    for s in students:
        att = s.get_attendance()
        if 'Slot A' in att:
            total_att_a += int(att['Slot A'] or 0)
            count_a += 1
    
    stats = {
        'total_students': len(students),
        'reports_generated': 152,
        'avg_attendance_a': int(total_att_a / count_a) if count_a > 0 else 0,
        'avg_attendance_b': 92
    }

    students_json = [
        {
            'reg_num': s.reg_num,
            'name': s.name,
            'course': s.course,
            'attendance_data': s.attendance_data,
            'marks_data': s.marks_data,
            'slot_info': s.slot_info,
            'registered_new_course': s.registered_new_course,
            'online_course': s.online_course,
            'event_participation': s.event_participation,
            'additional_description': s.additional_description,
            'photo_path': s.photo_path,
        }
        for s in students
    ]

    mentor_display = MENTOR_NAMES.get(session['user'], session['user'])
    return render_template('faculty.html', students=students, students_json=students_json, stats=stats, mentor_name=mentor_display)

@app.route('/generate_report')
def generate_report():
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    students = StudentDetail.query.filter_by(mentor_username=session['user']).all()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for student in students:
        slots = student.get_slots()
        att_data = student.get_attendance()
        marks_data = student.get_marks()
        mentor_name = MENTOR_NAMES.get(session['user'], session['user'])

        low_attendance_flag = False
        for sl in slots:
            try:
                if float(att_data.get(sl, 0)) < 80:
                    low_attendance_flag = True
                    break
            except ValueError:
                pass
        
        # --- SLIDE 1: ACADEMIC PERFORMANCE ---
        slide_layout = prs.slide_layouts[6] # Blank
        slide1 = prs.slides.add_slide(slide_layout)
        
        # Header Box with Logo-like text
        header_shape = slide1.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(12.9), Inches(1.2))
        tf = header_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "SIMATS ENGINEERING"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.name = "Times New Roman"
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        # Banner Table for Name/Reg/Mentor 1x6 for widescreen array
        banner_tbl = slide1.shapes.add_table(1, 6, Inches(0.2), Inches(1.4), Inches(12.9), Inches(0.5)).table
        banner_content = [
            "Mentee Name:", student.name or 'N/A', 
            "Reg.NO:", student.reg_num or 'N/A', 
            "Mentor name:", mentor_name
        ]
        col_widths = [Inches(1.8), Inches(3.0), Inches(1.3), Inches(2.0), Inches(1.8), Inches(3.0)]
        
        for i, text in enumerate(banner_content):
            banner_tbl.columns[i].width = col_widths[i]
            cell = banner_tbl.cell(0, i)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = "Times New Roman"
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

        # Student Photo with gray border
        left_img = Inches(0.5)
        top_img = Inches(2.5)
        width_img = Inches(3.0)
        height_img = Inches(3.5)
        
        # Draw a gray frame
        frame = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_img - Inches(0.1), top_img - Inches(0.1), width_img + Inches(0.2), height_img + Inches(0.2))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(230, 230, 230)
        frame.line.color.rgb = RGBColor(200, 200, 200)

        photo_added = False
        try:
            if student.photo_path:
                if student.photo_path.startswith('http'):
                    req = urllib.request.Request(student.photo_path, headers={'User-Agent': 'Mozilla/5.0'})
                    with urllib.request.urlopen(req) as response:
                        image_stream = BytesIO(response.read())
                    slide1.shapes.add_picture(image_stream, left_img, top_img, width=width_img, height=height_img)
                    photo_added = True
                elif os.path.exists(student.photo_path):
                    slide1.shapes.add_picture(student.photo_path, left_img, top_img, width=width_img, height=height_img)
                    photo_added = True
        except Exception as e:
            print(f"Error adding photo to PPT: {e}")

        if not photo_added:
            rect = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_img, top_img, width_img, height_img)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(100, 100, 100)
            text_frame = rect.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = "No Photo"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(16)
            p.font.name = "Times New Roman"
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Table for Marks
        rows = max(1 + len(slots), 2)
        cols = 4
        table_width = Inches(8.5)
        table_height = Inches(0.8 * rows)
        left_tbl = Inches(4.2)
        top_tbl = Inches(2.5)
        
        table = slide1.shapes.add_table(rows, cols, left_tbl, top_tbl, table_width, table_height).table
        
        # Header Row Styling
        h_labels = ["Course Slot", "Total Marks", "Marks Obtained", "Class Average Mark"]
        for i, h in enumerate(h_labels):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(112, 173, 71) # SIMATS Green
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.name = "Times New Roman"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

        row_data = []
        if not slots:
            row_data.append(["Test 1 (N/A)", "20", "0", "0"])
        else:
            for slot in slots:
                s_marks = marks_data.get(slot, {})
                test1_value = s_marks.get('test1', '0') or '0'
                avg_value = s_marks.get('avg', '0') or '0'
                total_marks_value = s_marks.get('total_marks', '20') or '20'
                row_data.append([f"Test 1 ({slot})", total_marks_value, test1_value, avg_value])
        
        for r_idx, r_vals in enumerate(row_data):
            for c_idx, val in enumerate(r_vals):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(226, 239, 218) # Light green
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(18)
                p.font.name = "Times New Roman"
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

        # --- SLIDE 2: MENTOR NOTES & ATTENDANCE ---
        slide2 = prs.slides.add_slide(slide_layout)
        
        # Header Box with Logo-like text
        header2 = slide2.shapes.add_textbox(Inches(0.2), Inches(0.2), Inches(12.9), Inches(1.2))
        tf2 = header2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "SIMATS ENGINEERING"
        p2.font.size = Pt(44)
        p2.font.bold = True
        p2.font.name = "Times New Roman"
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.alignment = PP_ALIGN.CENTER

        # Main Gray Content Box
        body_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8))
        body_box.fill.solid()
        body_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        body_box.line.color.rgb = RGBColor(200, 200, 200)

        tf_body = slide2.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.6)).text_frame
        tf_body.word_wrap = True
        
        # "Welcome to SIMATS ENGINEERING"
        p = tf_body.paragraphs[0]
        run = p.add_run()
        run.text = "Welcome to SIMATS ENGINEERING"
        run.font.bold = True
        run.font.size = Pt(22)
        run.font.name = "Times New Roman"
        run.font.color.rgb = RGBColor(0, 0, 0)
        add_highlight(run, '00FF00') # Green highlight
        
        p = tf_body.add_paragraph()
        run = p.add_run()
        run.text = "Dear Parent,"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Times New Roman"
        
        if slots:
            for slot in slots:
                att_val = att_data.get(slot, '0') or '0'
                s_marks = marks_data.get(slot, {}) or {}
                course_val = s_marks.get('course', 'N/A') or 'N/A'
                
                p = tf_body.add_paragraph()
                run = p.add_run()
                run.text = f"Attendance for {slot}: {course_val}: {att_val}%"
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.name = "Times New Roman"
                
                try:
                    is_low = float(att_val) < 80
                except ValueError:
                    is_low = False
                    
                if is_low:
                    add_highlight(run, 'FF0000') # Red highlighting
                else:
                    add_highlight(run, '00FF00') # Green highlighting
        else:
            p = tf_body.add_paragraph()
            run = p.add_run()
            run.text = "Attendance: N/A"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.name = "Times New Roman"
            add_highlight(run, '00FF00')
        
        p = tf_body.add_paragraph()
        p.space_before = Pt(18)
        run = p.add_run()
        run.text = f"{student.additional_description or 'I personally advised him to concentrate more on study and skill development... now he is currently attending an online course to improve his technical skills which is really appreciable...'}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Times New Roman"
        add_highlight(run, 'FFFF00') # Yellow highlighting
        
        p = tf_body.add_paragraph()
        run = p.add_run()
        run.text = f"New course: {student.registered_new_course or 'N/A'}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Times New Roman"
        add_highlight(run, 'FFFF00') # Yellow highlighting
        
        p = tf_body.add_paragraph()
        p.space_before = Pt(18)
        run = p.add_run()
        if student.event_participation:
            run.text = f"Your ward participated in: {student.event_participation} and gave his very best throughout the journey. His dedication, hard work, and sincere efforts are truly appreciable."
        else:
            run.text = "I have personally advised this student to actively participate in co-curricular and extracurricular events to enhance their overall personality, communication skills, and technical exposure. Regular participation in such events will greatly contribute to their professional development and help them stand out in their academics and career."
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Times New Roman"
        add_highlight(run, 'FFFF00') # Yellow highlighting

        # Administrative notices (Green Highlights)
        p = tf_body.add_paragraph()
        run = p.add_run()
        run.text = "All students are advised to pay their 2nd-year tuition fees on time through the Viana Portal."
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Times New Roman"
        add_highlight(run, '00FF00')

        p = tf_body.add_paragraph()
        run = p.add_run()
        run.text = "Additionally, kindly upload your recent passport-size photograph to your Viana profile at the earliest, if you have not already done so...."
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Times New Roman"
        add_highlight(run, '00FF00')

    report_path = 'Mentor_Dashboard_Report.pptx'
    prs.save(report_path)
    return send_file(report_path, as_attachment=True)

@app.route('/clear_student/<reg_num>', methods=['POST'])
def clear_student(reg_num):
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    student = StudentDetail.query.get(reg_num)
    if student:
        student.slot_info = '[]'
        student.attendance_data = '{}'
        student.marks_data = '{}'
        student.registered_new_course = None
        student.online_course = None
        student.event_participation = None
        student.additional_description = None
        delete_photo(student.photo_path)
        student.photo_path = None
        db.session.commit()
        flash(f'Report data for {student.name or reg_num} has been cleared.', 'success')
    else:
        flash('Student not found.', 'danger')
    
    return redirect(url_for('faculty_dashboard'))

@app.route('/remove_student/<reg_num>', methods=['POST'])
def remove_student(reg_num):
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    student = StudentDetail.query.get(reg_num)
    if student:
        delete_photo(student.photo_path)
        db.session.delete(student)
    
    # Delete the student's login account
    user = User.query.get(reg_num)
    if user:
        db.session.delete(user)
    
    db.session.commit()
    flash(f'Student {reg_num} has been permanently removed.', 'success')
    return redirect(url_for('faculty_dashboard'))

@app.route('/delete_all_reports', methods=['POST'])
def delete_all_reports():
    if 'user' not in session or session['role'] != 'faculty':
        return redirect(url_for('login'))
    
    # Only clear PPT/report fields — keep student accounts and rows intact
    students = StudentDetail.query.filter_by(mentor_username=session['user']).all()
    for student in students:
        delete_photo(student.photo_path)
        student.slot_info = '[]'
        student.attendance_data = '{}'
        student.marks_data = '{}'
        student.registered_new_course = None
        student.online_course = None
        student.event_participation = None
        student.additional_description = None
        student.photo_path = None

    db.session.commit()
    flash('All student report data has been cleared. Accounts are still active.', 'success')
    return redirect(url_for('faculty_dashboard'))

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
        seed_db()
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
